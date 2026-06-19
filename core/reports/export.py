"""
Multi-format report export: PDF and PPTX.

Both are optional — the libraries are only imported when the format is requested.
Install with: pip install weasyprint python-pptx
"""

from __future__ import annotations

import os
from pathlib import Path
from typing import Any

import structlog

logger = structlog.get_logger(__name__)


def get_report_formats() -> set[str]:
    raw = os.environ.get("REPORT_FORMAT", "json,html").strip()
    return {f.strip().lower() for f in raw.split(",") if f.strip()}


def export_pdf(report: dict[str, Any], output_path: Path) -> Path:
    try:
        from weasyprint import HTML  # type: ignore[import-untyped]
    except ImportError as exc:
        raise ImportError(
            "weasyprint is required for PDF export. "
            "Install with: pip install weasyprint"
        ) from exc

    from core.reports.html import build_html_report

    html_content = build_html_report(report)
    pdf_path = output_path.with_suffix(".pdf")
    HTML(string=html_content).write_pdf(str(pdf_path))
    logger.info("pdf_report_saved", path=str(pdf_path))
    return pdf_path


def export_pptx(report: dict[str, Any], output_path: Path) -> Path:
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
        from pptx.util import Inches  # type: ignore[import-untyped]
    except ImportError as exc:
        raise ImportError(
            "python-pptx is required for PPTX export. "
            "Install with: pip install python-pptx"
        ) from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"Argus — {report['cloud'].upper()} Waste Report"
    slide.placeholders[1].text = (
        f"Generated: {report['generated_at'][:10]}\n"
        f"Total estimated waste: ${report['total_estimated_waste_usd']:,.2f}/month\n"
        f"Findings: {report['findings_count']}"
    )

    # Executive summary slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    slide.placeholders[1].text = report.get("executive_summary", "")

    # Top findings slide(s) — max 8 per slide
    findings = report.get("findings", [])
    chunk_size = 8
    for i in range(0, len(findings), chunk_size):
        chunk = findings[i : i + chunk_size]
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = (
            f"Top Findings ({i + 1}–{i + len(chunk)} of {len(findings)})"
        )
        lines = []
        for f in chunk:
            name = f.get("name") or f["resource_id"]
            cost = f["estimated_monthly_cost"]
            priority = (f.get("priority") or "low").upper()
            lines.append(
                f"[{priority}] {name} — {f['resource_type']} — " f"${cost:,.2f}/mo"
            )
            lines.append(f"  → {f.get('recommendation', '')}")
        slide.placeholders[1].text = "\n".join(lines)

    # Scan metadata slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Scan Details"
    meta_lines = [
        f"Scan ID: {report['scan_id']}",
        f"Cloud: {report['cloud']}",
        f"Accounts: {', '.join(report.get('accounts_scanned', []))}",
    ]
    if report.get("agent_input_tokens"):
        meta_lines.append(
            f"AI tokens: {report['agent_input_tokens']:,} in / "
            f"{report['agent_output_tokens']:,} out"
        )
        meta_lines.append(
            f"Estimated AI cost: ${report.get('estimated_agent_cost_usd', 0):.4f}"
        )
    slide.placeholders[1].text = "\n".join(meta_lines)

    pptx_path = output_path.with_suffix(".pptx")
    prs.save(str(pptx_path))
    logger.info("pptx_report_saved", path=str(pptx_path))
    return pptx_path
